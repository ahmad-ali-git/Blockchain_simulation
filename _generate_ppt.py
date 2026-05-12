"""
Generates the Final Presentation as a .pptx file.
Output: FinalSubmission/Final_Presentation.pptx

Screenshots for GUI slides (add these files to embed real screenshots):
  FinalSubmission/screenshots/gui_intro.png      — intro animation overlay
  FinalSubmission/screenshots/gui_config.png     — configuration bento panel
  FinalSubmission/screenshots/gui_simulation.png — live simulation view
  FinalSubmission/screenshots/gui_results.png    — results dashboard
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = "FinalSubmission/Final_Presentation.pptx"
ASSETS = "FinalSubmission/ReportAssets"
SCREENSHOTS = "FinalSubmission/screenshots"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=GREY, align=PP_ALIGN.LEFT, bullet=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet:
            p.text = "•  " + line
        else:
            p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(8)
    return tb


def add_bullet_slide(title, bullets, image=None, image_left=8, image_width=4.7):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_title_bar(slide, title)
    width = 7 if image else 12
    add_textbox(slide, 0.6, 1.1, width, 5.8, bullets, size=20, bullet=True)
    if image:
        slide.shapes.add_picture(image, Inches(image_left), Inches(1.3),
                                 width=Inches(image_width))
    return slide


def add_screenshot_slide(title, image_path, caption="", fallback_bullets=None,
                         left_bullets=None):
    """Add a slide with a full-width screenshot or descriptive text fallback.

    If left_bullets is provided and the image exists, the layout is split:
    bullets on the left (~4.7"), screenshot on the right.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)

    img_exists = os.path.exists(image_path)

    if img_exists and left_bullets:
        # Split layout: bullets left, screenshot right
        add_textbox(slide, 0.5, 1.05, 4.5, 5.8, left_bullets, size=16, bullet=True)
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(0.95),
                                 width=Inches(7.8))
    elif img_exists:
        # Full-width screenshot
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0),
                                 width=Inches(12.33))
    else:
        # Fallback: descriptive text
        bullets = fallback_bullets or []
        add_textbox(slide, 0.5, 1.1, 12.3, 5.6, bullets, size=19, bullet=True)

    if caption:
        add_textbox(slide, 0.5, 6.85, 12.3, 0.5, caption,
                    size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    return slide


# ============================================================
# SLIDE 1 — title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

add_textbox(slide, 0.5, 1.4, 12.3, 0.7,
            "Final Year Project", size=22, bold=False,
            color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.2, 12.3, 1.5,
            "A Simulation-Based Framework for\nMitigating a 51% Attack on Blockchain Networks",
            size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.1, 12.3, 0.5,
            "Presented by Ahmad Ali  •  BC220425973", size=20,
            color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.65, 12.3, 0.5,
            "Group ID: F25PROJECT5043F", size=16,
            color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 6.3, 12.3, 0.5,
            "Supervisor: Dr. Fouzia Jumani  •  Virtual University of Pakistan",
            size=15, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — agenda
# ============================================================
add_bullet_slide(
    "Agenda",
    [
        "Background: what is a 51% attack",
        "Problem statement and project goal",
        "Proposed defense: Safe Mode Detection",
        "System architecture and design",
        "Implementation walkthrough",
        "GUI walkthrough (intro, configuration, live simulation, results, event logs)",
        "Live demonstration",
        "Results and analysis",
        "Limitations and future work",
        "Conclusion",
    ],
)

# ============================================================
# SLIDE 3 — background
# ============================================================
add_bullet_slide(
    "Background",
    [
        "Proof-of-Work blockchains agree on history using the longest chain rule",
        "Honest miners normally control the majority of the network's hashrate",
        "If one miner gains more than 50% of the hashrate, the rule starts to break down",
        "That miner can secretly build a private chain and publish it later to replace the public one",
        "This is what is called a 51% attack",
        "Several smaller PoW coins (Ethereum Classic, Bitcoin Gold, Verge) have already been hit",
    ],
)

# ============================================================
# SLIDE 4 — the 51% attack illustrated
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "How the Attack Works")
add_textbox(slide, 0.6, 1.1, 12, 0.6,
            "1. Attacker quietly forks the chain",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 1.8, 12, 0.6,
            "2. Mines blocks on the private chain in secret",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 2.5, 12, 0.6,
            "3. Inserts a double-spend transaction on the private fork",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 3.2, 12, 0.6,
            "4. Waits until the private chain is longer than the public one",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 3.9, 12, 0.6,
            "5. Publishes the private chain — the network adopts it because it is longer",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 4.6, 12, 0.6,
            "6. Original transactions are erased, attacker spends the same coins twice",
            size=20, bold=True, color=NAVY)
add_textbox(slide, 0.6, 5.6, 12, 1.2,
            "Once the attacker controls enough hashrate, this is not a matter of luck. "
            "Given enough time the longer chain will eventually be theirs.",
            size=18, italic=True, color=RED)

# ============================================================
# SLIDE 5 — problem statement
# ============================================================
add_bullet_slide(
    "Problem Statement",
    [
        "Smaller PoW networks have low total hashrate",
        "An attacker with modest mining hardware can rent enough cloud "
        "compute power to cross 50%",
        "Once they cross it, double-spending becomes possible",
        "Standard longest-chain rule cannot tell an honest reorganization "
        "from a fraudulent one",
        "We need a defense that does not rely on hashrate alone",
    ],
)

# ============================================================
# SLIDE 6 — project goal
# ============================================================
add_bullet_slide(
    "Project Goal",
    [
        "Build a Python simulation of a small PoW blockchain",
        "Make a malicious miner capable of carrying out the 51% attack",
        "Implement Safe Mode Detection and measure how it changes the result",
        "Provide a GUI that lets anyone run experiments interactively",
        "Compute hashrate from the number of mining nodes a miner owns "
        "(realistic, not a fixed number)",
    ],
)

# ============================================================
# SLIDE 7 — Safe Mode Detection
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Proposed Defense: Safe Mode Detection")

add_textbox(slide, 0.6, 1.1, 12, 0.5,
            "Safe Mode Detection has two parts that work together:",
            size=20, color=NAVY)

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(1.9), Inches(6), Inches(2.4))
box1.fill.solid()
box1.fill.fore_color.rgb = LIGHT
box1.line.color.rgb = NAVY
tf = box1.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "1. Long Private Chain (LPC)"
p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = ("Reject any block that would create six or more "
           "consecutive blocks by the same miner. This prevents "
           "any single party from dominating block production.")
p2.font.size = Pt(16); p2.font.color.rgb = GREY

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(6.8), Inches(1.9), Inches(6), Inches(2.4))
box2.fill.solid()
box2.fill.fore_color.rgb = LIGHT
box2.line.color.rgb = NAVY
tf = box2.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "2. UTXO Comparison"
p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = ("When a competing chain is published, compare the "
           "UTXO snapshots of both chains. Any difference is "
           "evidence of a double-spend, the chain is rejected.")
p2.font.size = Pt(16); p2.font.color.rgb = GREY

add_textbox(slide, 0.6, 4.7, 12, 1.5,
            "Together: even a > 50% attacker cannot replace the honest chain "
            "without either tripping the LPC streak limit or producing UTXO differences "
            "that the comparison step catches.",
            size=18, italic=True, color=GREEN)

# ============================================================
# SLIDE 8 — node-based hashrate (our approach only)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Dynamic Hashrate from Node Count")
add_textbox(slide, 0.6, 1.1, 12, 0.6,
            "Hashrate is not a fixed number — it is derived from the nodes a miner controls:",
            size=20, color=NAVY)
add_textbox(slide, 0.6, 2.0, 12, 0.5,
            "MaliciousMiner('Attacker', num_nodes=17)    # attacker owns 17 mining nodes",
            size=16, color=GREY)
add_textbox(slide, 0.6, 2.6, 12, 0.5,
            "HonestMiner('Honest_1', num_nodes=5)        # each honest miner owns 5 nodes",
            size=16, color=GREY)
add_textbox(slide, 0.6, 3.4, 12, 1.6,
            "Network computes at runtime:\n"
            "    hashrate = miner.num_nodes / total_network_nodes\n\n"
            "Example: 5 honest miners × 5 nodes = 25 honest nodes\n"
            "Attacker: 17 nodes  →  17 / 42 = 40.5% hashrate",
            size=18, color=NAVY)
add_textbox(slide, 0.6, 5.3, 12, 1.5,
            "This makes the simulation realistic: adding more mining hardware "
            "increases your share of the network proportionally, just like a real "
            "Proof-of-Work network.",
            size=18, italic=True, color=GREEN)

# ============================================================
# SLIDE 9 — architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "System Architecture (Three-Tier)")
slide.shapes.add_picture(f"{ASSETS}/diagram_architecture.png",
                         Inches(2.0), Inches(1.2), width=Inches(9.3))

# ============================================================
# SLIDE 10 — class overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Main Classes")
slide.shapes.add_picture(f"{ASSETS}/diagram_classes.png",
                         Inches(0.7), Inches(1.1), width=Inches(11.9))

# ============================================================
# SLIDE 11 — sequence baseline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Sequence: Baseline Attack")
_seq_base = f"{ASSETS}/sequence_baseline.png"
if os.path.exists(_seq_base):
    slide.shapes.add_picture(_seq_base, Inches(1.0), Inches(1.1), width=Inches(11.3))
else:
    add_textbox(slide, 0.6, 1.2, 12, 5.5, [
        "Honest miners broadcast new blocks → all miners accept them",
        "Attacker secretly forks: starts mining on a private chain",
        "Attacker adds double-spend tx to the private chain",
        "Public chain continues growing (honest miners only)",
        "When private chain length > public chain length: attacker publishes",
        "Network adopts the longer (attacker) chain — honest blocks orphaned",
        "Double-spend is now in the canonical chain — attack succeeds",
    ], size=19, bullet=True, color=NAVY)

# ============================================================
# SLIDE 12 — sequence protected
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Sequence: Protected Run")
_seq_prot = f"{ASSETS}/sequence_protected.png"
if os.path.exists(_seq_prot):
    slide.shapes.add_picture(_seq_prot, Inches(1.0), Inches(1.1), width=Inches(11.3))
else:
    add_textbox(slide, 0.6, 1.2, 12, 5.5, [
        "LPC Defense monitors every proposed block",
        "If the proposer already mined the last N consecutive blocks: block rejected",
        "Attacker builds private chain (same as baseline) while public chain grows",
        "When attacker tries to publish: LPC scans the proposed chain",
        "If 6+ consecutive blocks by same miner found: chain REJECTED",
        "If not: UTXO snapshot compared — any double-spend difference: REJECTED",
        "Honest chain remains canonical — attack blocked without needing majority hashrate",
    ], size=19, bullet=True, color=NAVY)

# ============================================================
# SLIDE 13 — implementation tools
# ============================================================
add_bullet_slide(
    "Tools and Libraries",
    [
        "Python 3.10 — programming language",
        "Streamlit — interactive web GUI built in pure Python",
        "Matplotlib + Plotly — charts (static for report, interactive for GUI)",
        "Pandas — tabulating results and CSV export",
        "hashlib (stdlib) — SHA-256 for block hashing",
        "Visual Studio Code — development environment",
    ],
)

# ============================================================
# SLIDE 14 — GUI: intro animation
# ============================================================
add_screenshot_slide(
    "GUI: Intro Animation",
    image_path=f"{SCREENSHOTS}/gui_intro.png",
    caption="Full-screen intro with animated network graph and project title — press SKIP to enter the app",
    fallback_bullets=[
        "Full-screen dark overlay on first load",
        "Animated network graph: nodes connected by blue (honest) and red (attacker) edges",
        "Project title centered: '51% ATTACK SIMULATION'",
        "Subtitle: 'A Simulation-Based Framework for Mitigating a 51% Attack on Blockchain Networks'",
        "Student info: Ahmad Ali · BC220425973 · F25PROJECT5043F",
        "SKIP button — takes the user directly to the configuration panel",
    ],
)

# ============================================================
# SLIDE 15 — GUI: configuration panel
# ============================================================
add_screenshot_slide(
    "GUI: Configuration Panel",
    image_path=f"{SCREENSHOTS}/gui_config.png",
    caption="3-column bento grid: honest miners (with node-dot grid), attacker (live hashrate ring), parameters",
    fallback_bullets=[
        "Animated header: scrolling chain of blocks with hash labels",
        "3-column bento grid for all simulation parameters",
        "Honest miners card: stepper controls for miners and nodes/miner, visual node-dot grid",
        "Attacker card: node count stepper, live % ring gauge, WARNING badge approaching majority",
        "Parameters card: total blocks, LPC limit, PoW difficulty, run mode selector",
        "Network ring topology below: blue dots = honest nodes, red dots = attacker nodes",
        "Live hashrate stats: Honest Network 59.5% vs Attacker 40.5%",
    ],
    left_bullets=[
        "Stepper controls (not sliders) for miners/nodes",
        "Visual node-dot grid per miner",
        "Live % ring gauge for attacker hashrate",
        "WARNING badge when attacker nears 50%",
        "Network ring topology diagram",
        "Run mode: Compare Both / Baseline Only / Protected Only",
    ],
)

# ============================================================
# SLIDE 16 — GUI: live simulation
# ============================================================
add_screenshot_slide(
    "GUI: Live Simulation View",
    image_path=f"{SCREENSHOTS}/gui_simulation.png",
    caption="Baseline (top) and Protected (bottom) sections run inline — chain growth chart + color-coded event log",
    fallback_bullets=[
        "Single-page flow: simulation renders inline, no tab switching",
        "Each section has a progress bar, chain growth chart, and event log side by side",
        "Chain growth chart: cyan = honest chain, red dashed = attacker private chain",
        "Green log entries — block mined by honest miner",
        "Red log entries   — attacker activity (mining secretly / published chain)",
        "Blue log entries  — LPC defense rejections",
        "Baseline section (top) runs first, Protected (bottom) runs after",
    ],
    left_bullets=[
        "Inline rendering — no tab switching",
        "Chain growth: honest (cyan) vs attacker (red dashed)",
        "Green = mined / Red = attack / Blue = defense",
        "Baseline runs first, Protected runs after",
        "Progress bar per section",
    ],
)

# ============================================================
# SLIDE 16b — GUI: results dashboard
# ============================================================
add_screenshot_slide(
    "GUI: Results Dashboard",
    image_path=f"{SCREENSHOTS}/gui_results.png",
    caption="SUCCEEDED (red) vs BLOCKED (green) outcome cards | 4-line chain growth | miner bar chart | chain viz",
    fallback_bullets=[
        "Outcome status cards: SUCCEEDED in red (attacker won) / BLOCKED in green (defense held)",
        "Comparison table: attack outcome, final chain length, blocks rejected, attacker hashrate",
        "4-line chain growth chart: Baseline-Honest, Baseline-Attacker, Protected-Honest, Protected-Attacker",
        "Miner statistics: horizontal bar chart with blocks mined per miner (Attacker in red, honest in cyan)",
        "Final chain visualization: color-coded block-by-block strip showing who mined each block",
    ],
    left_bullets=[
        "SUCCEEDED (red) vs BLOCKED (green) status cards",
        "Comparison table for both scenarios",
        "4-line chain growth chart",
        "Horizontal bar: blocks mined per miner",
        "Color-coded chain block strip",
    ],
)

# ============================================================
# SLIDE 16c — GUI: miner stats & event logs
# ============================================================
add_screenshot_slide(
    "GUI: Miner Stats & Event Logs",
    image_path=f"{SCREENSHOTS}/gui_stats.png",
    caption="Miner table (nodes, hashrate, blocks mined) | Protected chain viz | color-coded event logs + CSV download",
    fallback_bullets=[
        "Miner stats table: Miner ID, Nodes, Hashrate %, Blocks Mined — verified correct for all miners",
        "Protected chain visualization strip (fewer attacker blocks — defense worked)",
        "Full event logs for both Baseline and Protected runs (scrollable, color-coded)",
        "Download Baseline Log (CSV) and Download Protected Log (CSV) buttons",
    ],
    left_bullets=[
        "Miner table: nodes, hashrate, blocks mined",
        "Protected chain viz strip",
        "Full color-coded event logs",
        "CSV download for both runs",
    ],
)

# ============================================================
# SLIDE 17 — live demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Live Demo")
add_textbox(slide, 0.6, 1.5, 12, 1,
            "Run the GUI:",
            size=24, bold=True, color=NAVY)
add_textbox(slide, 0.6, 2.4, 12, 0.7,
            "    streamlit run app.py",
            size=22, color=GREY)
add_textbox(slide, 0.6, 3.5, 12, 3,
            "Demonstrate:\n"
            "1. Adjust attacker nodes — watch the hashrate readout cross 50%\n"
            "2. Run Compare mode — both scenarios execute side by side\n"
            "3. Baseline result: attacker SUCCEEDS (chain replaced)\n"
            "4. Protected result: LPC + UTXO check BLOCKS the attack\n"
            "5. Inspect the chain visualizer and event log details",
            size=20, color=NAVY, bullet=False)

# ============================================================
# SLIDE 18 — results: main chart
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Result 1: Attack Success Rate")
slide.shapes.add_picture(f"{ASSETS}/chart_attack_success.png",
                         Inches(2.5), Inches(1.1), width=Inches(8.3))
add_textbox(slide, 0.6, 6.4, 12, 0.9,
            "Baseline 47% success vs Protected 0% — Safe Mode catches every attack at 54.5% hashrate",
            size=16, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 19 — results: hashrate sweep
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Result 2: Sweep Across Attacker Hashrates")
slide.shapes.add_picture(f"{ASSETS}/chart_hashrate_vs_success.png",
                         Inches(2.5), Inches(1.1), width=Inches(8.3))
add_textbox(slide, 0.6, 6.4, 12, 0.9,
            "Even at 70%+ attacker hashrate, the protected line stays at 0%",
            size=16, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 20 — chain growth single run
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Result 3: Single-Run Chain Growth")
slide.shapes.add_picture(f"{ASSETS}/chart_chain_growth.png",
                         Inches(2.5), Inches(1.1), width=Inches(8.3))
add_textbox(slide, 0.6, 6.4, 12, 0.9,
            "Honest vs attacker private chain over time — until the fork point only the public chain grows",
            size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 21 — limitations
# ============================================================
add_bullet_slide(
    "Limitations",
    [
        "Mining is modeled probabilistically (weighted random selection), "
        "not real hash competition — standard for academic simulators",
        "No network latency or partition modeling",
        "Only the secret-fork-and-publish attack pattern is implemented",
        "Difficulty is fixed; real Bitcoin adjusts it every 2016 blocks",
        "These are deliberate scope limits, documented in the report",
    ],
)

# ============================================================
# SLIDE 22 — future work
# ============================================================
add_bullet_slide(
    "Future Work",
    [
        "Add network latency and partition modeling for realism",
        "Implement and compare other defenses (checkpointing, finality gadgets)",
        "Add other attack patterns: selfish mining, eclipse attacks",
        "Scale up to many more miners and longer chains",
        "Wrap the engine into a reusable Python package others can extend",
    ],
)

# ============================================================
# SLIDE 23 — conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Conclusion")
add_textbox(slide, 0.6, 1.2, 12, 5,
            [
                "Built a working Python simulation of a Proof-of-Work blockchain",
                "Implemented a malicious miner that carries out 51% attacks",
                "Implemented Safe Mode Detection: LPC + UTXO comparison",
                "Hashrate is derived dynamically from node counts, as the supervisor required",
                "Streamlit GUI lets anyone run experiments and watch the attack live",
                "Results: 47% baseline attack success dropped to 0% with the defense",
            ],
            size=20, bullet=True, color=NAVY)
add_textbox(slide, 0.6, 6.4, 12, 0.7,
            "The defense works under the conditions tested. Limitations and future work are documented in the report.",
            size=15, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 24 — thank you
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
add_textbox(slide, 0.5, 2.7, 12.3, 1,
            "Thank You", size=72, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.0, 12.3, 0.7,
            "Questions and feedback welcome", size=24,
            color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.5, 12.3, 0.5,
            "Ahmad Ali  •  BC220425973  •  F25PROJECT5043F", size=16,
            color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"Presentation written to {OUT}")
print(f"Total slides: {len(prs.slides)}")
